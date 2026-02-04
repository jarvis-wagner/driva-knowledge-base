#!/usr/bin/env python3
"""
Baixa apresentações do Google Drive e extrai texto para .md

Uso:
    python scripts/process_presentations.py
"""
from google.oauth2.credentials import Credentials
from googleapiclient.discovery import build
from googleapiclient.http import MediaIoBaseDownload
from pptx import Presentation
import json
import io
import os
from pathlib import Path

BASE_DIR = Path(__file__).parent.parent
TOKEN_FILE = BASE_DIR.parent / ".google_token.json"
OUTPUT_DIR = BASE_DIR / "produtos" / "apresentacoes"
TEMP_DIR = BASE_DIR / "source_files" / "presentations"

# Pasta raiz no Drive
FOLDER_ID = "14hZDQzNs5gq6b_u5kM0NH77wTZvKWBTY"


def get_credentials():
    with open(TOKEN_FILE, 'r') as f:
        token_data = json.load(f)
    return Credentials(
        token=token_data['token'],
        refresh_token=token_data['refresh_token'],
        token_uri=token_data['token_uri'],
        client_id=token_data['client_id'],
        client_secret=token_data['client_secret'],
        scopes=token_data['scopes']
    )


def list_presentations(service, folder_id, path=""):
    """Lista apresentações recursivamente"""
    presentations = []
    
    results = service.files().list(
        q=f"'{folder_id}' in parents and trashed=false",
        fields="files(id, name, mimeType)",
        supportsAllDrives=True,
        includeItemsFromAllDrives=True,
        pageSize=100
    ).execute()
    
    for f in results.get('files', []):
        current_path = f"{path}/{f['name']}" if path else f['name']
        
        if f['mimeType'] == 'application/vnd.google-apps.folder':
            # Recursão para subpastas
            presentations.extend(list_presentations(service, f['id'], current_path))
        elif f['mimeType'] == 'application/vnd.google-apps.presentation':
            presentations.append({
                'id': f['id'],
                'name': f['name'],
                'path': current_path,
                'mimeType': f['mimeType']
            })
    
    return presentations


def download_as_pptx(service, file_id, output_path):
    """Baixa Google Slides como PPTX"""
    request = service.files().export_media(
        fileId=file_id,
        mimeType='application/vnd.openxmlformats-officedocument.presentationml.presentation'
    )
    
    fh = io.BytesIO()
    downloader = MediaIoBaseDownload(fh, request)
    done = False
    while not done:
        status, done = downloader.next_chunk()
    
    fh.seek(0)
    output_path.parent.mkdir(parents=True, exist_ok=True)
    output_path.write_bytes(fh.read())
    return output_path


def extract_text_from_pptx(pptx_path):
    """Extrai texto de um arquivo PPTX"""
    prs = Presentation(pptx_path)
    slides_content = []
    
    for i, slide in enumerate(prs.slides, 1):
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_text.append(shape.text.strip())
        
        if slide_text:
            slides_content.append({
                'slide': i,
                'content': slide_text
            })
    
    return slides_content


def generate_markdown(name, path, slides):
    """Gera conteúdo markdown a partir dos slides"""
    md = f"# {name}\n\n"
    md += f"> Fonte: {path}\n\n"
    md += "---\n\n"
    
    for slide in slides:
        md += f"## Slide {slide['slide']}\n\n"
        for text in slide['content']:
            # Limpar e formatar
            text = text.replace('\n', '\n\n')
            md += f"{text}\n\n"
        md += "---\n\n"
    
    return md


def main():
    print("🔐 Carregando credenciais...")
    creds = get_credentials()
    service = build('drive', 'v3', credentials=creds)
    
    print("🔍 Listando apresentações...")
    presentations = list_presentations(service, FOLDER_ID)
    print(f"   Encontradas: {len(presentations)} apresentações\n")
    
    # Criar diretórios
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    TEMP_DIR.mkdir(parents=True, exist_ok=True)
    
    processed = 0
    for pres in presentations:
        print(f"📄 {pres['name'][:50]}...")
        
        try:
            # Baixar como PPTX
            safe_name = "".join(c if c.isalnum() or c in "._- " else "_" for c in pres['name'])
            pptx_path = TEMP_DIR / f"{safe_name}.pptx"
            
            download_as_pptx(service, pres['id'], pptx_path)
            print(f"   ✅ Baixado")
            
            # Extrair texto
            slides = extract_text_from_pptx(pptx_path)
            print(f"   📝 {len(slides)} slides com conteúdo")
            
            # Gerar markdown
            md_content = generate_markdown(pres['name'], pres['path'], slides)
            
            # Salvar
            md_filename = f"{safe_name}.md"
            md_path = OUTPUT_DIR / md_filename
            md_path.write_text(md_content, encoding='utf-8')
            print(f"   💾 Salvo: {md_filename}")
            
            processed += 1
            
        except Exception as e:
            print(f"   ❌ Erro: {e}")
        
        print()
    
    print(f"✅ Processadas: {processed}/{len(presentations)} apresentações")
    print(f"📁 Output: {OUTPUT_DIR}")


if __name__ == "__main__":
    main()
